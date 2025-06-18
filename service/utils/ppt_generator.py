"""
PowerPoint generation utilities for creating startup presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional
import os
import tempfile
import json
import logging

logger = logging.getLogger(__name__)

class PowerPointGenerator:
    """
    Service for generating PowerPoint presentations for startups.
    """
    
    def __init__(self):
        self.slide_layouts = {
            'title': 0,
            'content': 1,
            'section_header': 2,
            'two_content': 3,
            'comparison': 4,
            'title_only': 5,
            'blank': 6,
            'content_with_caption': 7,
            'picture_with_caption': 8
        }
    
    def create_pitch_deck(self, content_structure: dict, user_info: dict) -> str:
        """Create a pitch deck presentation with the provided content"""
        try:
            # Create presentation
            prs = Presentation()
            
            # Add slides based on content
            self._add_title_slide(prs, content_structure.get('title', {}))
            self._add_problem_slide(prs, content_structure.get('problem', {}))
            self._add_solution_slide(prs, content_structure.get('solution', {}))
            self._add_market_slide(prs, content_structure.get('market', {}))
            self._add_business_model_slide(prs, content_structure.get('business_model', {}))
            self._add_competition_slide(prs, content_structure.get('competition', {}))
            self._add_team_slide(prs, content_structure.get('team', {}))
            self._add_financials_slide(prs, content_structure.get('financials', {}))
            self._add_funding_slide(prs, content_structure.get('funding', {}))
            self._add_next_steps_slide(prs, content_structure.get('next_steps', {}))
            
            # Save to temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            prs.save(temp_file.name)
            temp_file.close()
            
            logger.info(f"Pitch deck generated successfully: {temp_file.name}")
            return temp_file.name
            
        except Exception as e:
            logger.error(f"Error creating pitch deck: {str(e)}")
            raise
    
    def create_business_plan_presentation(self, content_data: Dict[str, Any]) -> str:
        """
        Create a business plan presentation.
        
        Args:
            content_data: Dictionary containing all the content for the business plan
            
        Returns:
            str: Path to the generated PowerPoint file
        """
        try:
            # Create presentation
            prs = Presentation()
            
            # Add business plan slides
            self._add_title_slide(prs, content_data.get('title', {}))
            self._add_executive_summary_slide(prs, content_data.get('executive_summary', {}))
            self._add_company_description_slide(prs, content_data.get('company_description', {}))
            self._add_market_analysis_slide(prs, content_data.get('market_analysis', {}))
            self._add_organization_slide(prs, content_data.get('organization', {}))
            self._add_service_description_slide(prs, content_data.get('service_description', {}))
            self._add_marketing_sales_slide(prs, content_data.get('marketing_sales', {}))
            self._add_funding_request_slide(prs, content_data.get('funding_request', {}))
            self._add_financial_projections_slide(prs, content_data.get('financial_projections', {}))
            self._add_appendix_slide(prs, content_data.get('appendix', {}))
            
            # Save to temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            prs.save(temp_file.name)
            temp_file.close()
            
            logger.info(f"Business plan presentation generated successfully: {temp_file.name}")
            return temp_file.name
            
        except Exception as e:
            logger.error(f"Error creating business plan presentation: {str(e)}")
            raise
    
    def _add_title_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add title slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['title']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content.get('title', 'Startup Presentation')
        subtitle.text = content.get('subtitle', 'Building the Future')
        
        # Style the title
        title_format = title.text_frame.paragraphs[0].font
        title_format.size = Pt(44)
        title_format.bold = True
        title_format.color.rgb = RGBColor(68, 114, 196)
    
    def _add_problem_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add problem slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Problem"
        
        # Add problem description
        tf = content_placeholder.text_frame
        tf.text = content.get('description', 'Define the problem you are solving')
        
        # Add bullet points for pain points
        pain_points = content.get('pain_points', [])
        for pain_point in pain_points:
            p = tf.add_paragraph()
            p.text = pain_point
            p.level = 1
    
    def _add_solution_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add solution slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Solution"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('description', 'Our innovative solution')
        
        # Add key features
        features = content.get('features', [])
        for feature in features:
            p = tf.add_paragraph()
            p.text = feature
            p.level = 1
    
    def _add_market_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add market slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Market Opportunity"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('market_size', 'Large and growing market')
        
        # Add market details
        market_details = content.get('details', [])
        for detail in market_details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_business_model_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add business model slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Business Model"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('revenue_model', 'How we make money')
        
        # Add revenue streams
        revenue_streams = content.get('revenue_streams', [])
        for stream in revenue_streams:
            p = tf.add_paragraph()
            p.text = stream
            p.level = 1
    
    def _add_competition_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add competition slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Competition"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('competitive_advantage', 'Our competitive advantage')
        
        # Add competitors
        competitors = content.get('competitors', [])
        for competitor in competitors:
            p = tf.add_paragraph()
            p.text = competitor
            p.level = 1
    
    def _add_team_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add team slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Team"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('overview', 'Meet our amazing team')
        
        # Add team members
        team_members = content.get('members', [])
        for member in team_members:
            p = tf.add_paragraph()
            p.text = f"{member.get('name', '')} - {member.get('role', '')}"
            p.level = 1
    
    def _add_financials_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add financials slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Financial Projections"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('overview', 'Financial highlights and projections')
        
        # Add financial metrics
        metrics = content.get('metrics', [])
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = metric
            p.level = 1
    
    def _add_funding_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add funding slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Funding Request"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('amount', 'Investment opportunity')
        
        # Add use of funds
        use_of_funds = content.get('use_of_funds', [])
        for use in use_of_funds:
            p = tf.add_paragraph()
            p.text = use
            p.level = 1
    
    def _add_next_steps_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add next steps slide to presentation"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Next Steps"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('overview', 'Our roadmap forward')
        
        # Add steps
        steps = content.get('steps', [])
        for step in steps:
            p = tf.add_paragraph()
            p.text = step
            p.level = 1
    
    # Business Plan specific slides
    def _add_executive_summary_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add executive summary slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Executive Summary"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('summary', 'Executive summary of the business')
        
        # Add key points
        key_points = content.get('key_points', [])
        for point in key_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
    
    def _add_company_description_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add company description slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Company Description"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('description', 'About our company')
        
        # Add company details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_market_analysis_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add market analysis slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Market Analysis"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('analysis', 'Comprehensive market analysis')
        
        # Add analysis points
        points = content.get('points', [])
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
    
    def _add_organization_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add organization slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Organization & Management"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('structure', 'Organizational structure')
        
        # Add structure details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_service_description_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add service description slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Service or Product Line"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('description', 'Our products and services')
        
        # Add service details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_marketing_sales_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add marketing and sales slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Marketing & Sales"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('strategy', 'Marketing and sales strategy')
        
        # Add strategy details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_funding_request_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add funding request slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Funding Request"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('request', 'Funding requirements')
        
        # Add funding details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_financial_projections_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add financial projections slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Financial Projections"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('projections', 'Financial forecasts')
        
        # Add projection details
        details = content.get('details', [])
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
    
    def _add_appendix_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add appendix slide"""
        slide_layout = prs.slide_layouts[self.slide_layouts['content']]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Appendix"
        
        tf = content_placeholder.text_frame
        tf.text = content.get('content', 'Additional information')
        
        # Add appendix items
        items = content.get('items', [])
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
